"""Generate an EDITABLE PowerPoint of the team-construction slide.

Every element is a native shape (rectangle / connector / text box) — no images —
so the deck owner can retype labels, drag boxes and recolour in PowerPoint.

    ~/.pyenv/versions/3.11.13/envs/autoAI/bin/python brainstorm/_team_construction_pptx.py

Geometry note: the design was authored at 1280x720 CSS px (see the sibling
2026-08-18-team-construction-slide.html) and a 16:9 slide is exactly 13.333in x
7.5in = 1280x720 at 96 DPI, so design px map 1:1 onto slide px. CSS px -> pt is
x0.75. Both conversions live in `px()` / `pt_of()` below, and every coordinate
in this file is a design px lifted straight from the HTML.
"""
from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ── palette (AmEx) ──────────────────────────────────────────────────────────
NAVY = RGBColor(0x00, 0x17, 0x5A)
BLUE = RGBColor(0x00, 0x6F, 0xCF)
GREEN = RGBColor(0x0B, 0x8A, 0x5F)
GREEN_D = RGBColor(0x0A, 0x5C, 0x42)
RED = RGBColor(0xA0, 0x35, 0x30)
GRAY_TX = RGBColor(0x4A, 0x55, 0x68)
GRAY_MID = RGBColor(0x6B, 0x7A, 0x90)
GRAY_LN = RGBColor(0xE2, 0xE8, 0xF0)
GRAY_BG = RGBColor(0xF4, 0xF6, 0xF9)
BOX_BG = RGBColor(0xF4, 0xF8, 0xFD)
BOX_LN = RGBColor(0xCF, 0xE0, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE_BLUE = RGBColor(0x9C, 0xC4, 0xE8)

FONT = "Arial"          # cross-platform stand-in for Helvetica Neue
MONO = "Consolas"       # ships with Office on both Windows and macOS
INK = RGBColor(0x1C, 0x24, 0x33)      # dark code-block ground
CODE_TX = RGBColor(0xD6, 0xDE, 0xEA)  # code-block body text
CODE_NUM = RGBColor(0x7F, 0xD1, 0xA8) # ground-truth numbers
CODE_STR = RGBColor(0x9C, 0xC4, 0xE8) # string literals
EMU_PER_PX = 9525       # 914400 EMU/in / 96 px/in


def px(v: float) -> Emu:
    return Emu(int(round(v * EMU_PER_PX)))


def pt_of(css_px: float) -> Pt:
    return Pt(css_px * 0.75)


# ── primitives ──────────────────────────────────────────────────────────────
def rect(shapes, x, y, w, h, fill=None, line=None, line_w=0.75):
    sh = shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.text_frame.word_wrap = False
    return sh


def _style_runs(para, text, size, bold, color, font=FONT):
    """Split `text` on ** ** and emit bold/normal runs."""
    for i, chunk in enumerate(text.split("**")):
        if not chunk:
            continue
        r = para.add_run()
        r.text = chunk
        r.font.size = pt_of(size)
        r.font.name = font
        r.font.bold = bold or (i % 2 == 1)
        r.font.color.rgb = color


def text(shapes, x, y, w, h, body, size=9, bold=False, color=NAVY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, wrap=False,
         spacing=1.0, pad_l=0.0):
    """A text box. `body` may be a str or a list of (str, size, bold, color)."""
    tb = shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = px(pad_l)
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    lines = body if isinstance(body, list) else [(body, size, bold, color)]
    for i, spec in enumerate(lines):
        s, sz, bd, col = spec if isinstance(spec, tuple) else (spec, size, bold, color)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if i:
            p.space_before = Pt(0)
        _style_runs(p, s, sz, bd, col)
    return tb


def label_in(shape, body, size=8.5, bold=False, color=NAVY,
             align=PP_ALIGN.CENTER, pad_l=0.0):
    """Put centred text inside an existing shape."""
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = px(pad_l)
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = body if isinstance(body, list) else [body]
    for i, spec in enumerate(lines):
        s, sz, bd, col = spec if isinstance(spec, tuple) else (spec, size, bold, color)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 0.92
        _style_runs(p, s, sz, bd, col)
    return shape


def arrow(shapes, x1, y1, x2, y2, color=GRAY_MID, width=1.0, head=True):
    cn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    if head:
        ln = cn.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"),
                              {"type": "triangle", "w": "sm", "len": "sm"})
        ln.append(tail)
    return cn


def hline(shapes, x, y, w, color=GRAY_LN, width=0.75):
    cn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x), px(y), px(x + w), px(y))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    return cn


def chip(shapes, x, y, w, h, name, note=None, note_x=None):
    """Pale specialist box; optional grey note to its right, inside the box."""
    rect(shapes, x, y, w, h, BOX_BG, BOX_LN)
    text(shapes, x + 7, y, 120, h, name, 8.5, False, NAVY)
    if note:
        text(shapes, note_x if note_x is not None else x + 74, y, 200, h,
             note, 8.5, False, GRAY_MID)


def code(shapes, x, y, w, h, lines, size=9, color=WHITE, pad_l=8):
    """Monospaced block, one paragraph per line. `lines` items may be
    (str, color) to recolour a single line."""
    tb = shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = px(pad_l)
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, spec in enumerate(lines):
        s, col = spec if isinstance(spec, tuple) else (spec, color)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.22
        _style_runs(p, s, size, False, col, font=MONO)
    return tb


def new_slide(prs):
    """Blank slide carrying the deck's navy top rule."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide.shapes, 0, 0, 1280, 5, NAVY)
    return slide


def head(sh, title):
    text(sh, 44, 22, 1100, 42, title, 32, True, NAVY, anchor=MSO_ANCHOR.MIDDLE)
    hline(sh, 44, 72, 1192, RGBColor(0xC7, 0xCF, 0xDA), 1.0)


# ── slide 1 — team construction ─────────────────────────────────────────────
def slide_team_construction(prs):
    slide = new_slide(prs)
    sh = slide.shapes
    head(sh, "Team Construction | three decisions per turn")

    BANDS = [84, 286, 488]        # band top edges
    DX = 278                      # diagram origin x
    DIAG_DY = 29                  # diagram origin offset within band

    # ── left column text ────────────────────────────────────────────────
    left_col = [
        ("01", "Dispatch shape",
         "**Parallel by default.** Add a round only when one answer must "
         "exist before the next question can be asked."),
        ("02", "Team size",
         "**Minimal fetches. Maximal finds.** Open, senior-level questions "
         "deliberately invert the minimum-set rule."),
        ("03", "General specialist review",
         "**Server-enforced, not model-chosen.** On every ≥2-specialist "
         "turn, the team is checked for coherence before the answer ships."),
    ]
    for (num, title, principle), by in zip(left_col, BANDS):
        text(sh, 44, by + 26, 208, 14, num, 11, True, BLUE)
        text(sh, 44, by + 42, 208, 52, title, 20, True, NAVY,
             anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.05)
        text(sh, 44, by + 100, 208, 78, principle, 12.5, False, GRAY_TX,
             anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.25)

    # ── examples (right column) ─────────────────────────────────────────
    examples = [
        ["**“any spending spikes? what are the drivers?”**|q",
         "Causal — the spike window must exist before drivers can be read "
         "against it.|p",
         "**Naive parallel:** spend returns a spike in **2025-05**, modeling "
         "returns drivers from **2024-09**. Two halves that don’t connect.|r"],
        ["**“what’s the FICO?” → 1 specialist**|q",
         "Named metric, owning column, one fetch.|p",
         "**“any model opportunities?” → 4+**|q",
         "Names no metric and no table. Its value is in what the team "
         "**finds**, so each specialist gets a different hypothesis to test "
         "— and **a checked non-finding is a result**.|p"],
        ["**The 2025-05 vs 2024-09 mismatch, caught**|q",
         "Review sees both payloads, spots the windows don’t line up, and "
         "re-runs **modeling alone** anchored to **2025-05** — not a whole "
         "re-plan.|p",
         "Its phase-1 knowledge points are dropped first, so the corrected run "
         "doesn’t leave **two contradictory anchors** in memory.|p"],
    ]
    for lines, by in zip(examples, BANDS):
        rect(sh, 920, by + 14, 316, 162, GRAY_BG)
        rect(sh, 920, by + 14, 4, 162, BLUE)          # left accent bar
        text(sh, 936, by + 26, 292, 12, "E X A M P L E", 9.5, True, BLUE)

        y = by + 44
        for spec in lines:
            body, kind = spec.rsplit("|", 1)
            size = 12.5 if kind == "q" else 11.8
            col = NAVY if kind == "q" else (RED if kind == "r" else GRAY_TX)
            # rough line-count for vertical advance
            n = max(1, int(len(body.replace("**", "")) / 46) + 1)
            h = n * (16 if kind == "q" else 15)
            text(sh, 936, y, 292, h, body, size, kind == "q", col,
                 anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.18)
            y += h + (5 if kind != "q" else 4)

    # ═══ BAND 01 — dispatch shape ═══════════════════════════════════════
    by = BANDS[0] + DIAG_DY

    def D(x=0.0, y=0.0):
        return DX + x, by + y

    # PARALLEL
    text(sh, *D(0, -2), 200, 12, "P A R A L L E L  —  D E F A U L T", 10, True, BLUE)
    label_in(rect(sh, *D(0, 46), 40, 22, NAVY), "ORCH", 8, True, WHITE)
    for ey in (36, 57, 78):
        arrow(sh, *D(40, 57), *D(58, ey))
    for ry, nm in ((26, "spend_payments"), (48, "modeling"), (70, "crossbu")):
        chip(sh, DX + 66, by + ry, 124, 17, nm)
    text(sh, *D(0, 92), 200, 12, "independent sub-questions", 9.5, False, GRAY_MID)
    text(sh, *D(0, 105), 200, 12, "one round", 9.5, True, NAVY)

    # COLLAPSE
    text(sh, *D(208, -2), 200, 12, "C O L L A P S E", 10, True, BLUE)
    label_in(rect(sh, *D(208, 46), 40, 22, NAVY), "ORCH", 8, True, WHITE)
    arrow(sh, *D(248, 57), *D(268, 57))
    label_in(rect(sh, *D(276, 40), 122, 34, BOX_BG, BOX_LN),
             [("modeling", 8.5, True, NAVY),
              ("finds the window itself", 8, False, GRAY_MID)])
    text(sh, *D(276, 22), 122, 14, "↺  cross-queries", 8, True, BLUE,
         align=PP_ALIGN.CENTER)
    text(sh, *D(208, 92), 200, 12, "one specialist self-anchors", 9.5, False, GRAY_MID)
    text(sh, *D(208, 105), 200, 12, "one round, no handoff", 9.5, True, NAVY)

    # SEQUENTIAL
    text(sh, *D(416, -2), 200, 12, "S E Q U E N T I A L", 10, True, BLUE)
    label_in(rect(sh, *D(416, 46), 80, 22, BOX_BG, BOX_LN), "spend_payments", 8.5)
    arrow(sh, *D(496, 57), *D(520, 57), BLUE, 1.25)
    text(sh, *D(490, 40), 36, 12, "2025-05", 8, True, BLUE, align=PP_ALIGN.CENTER)
    label_in(rect(sh, *D(526, 46), 80, 22, BOX_BG, BOX_LN), "modeling", 8.5)
    text(sh, *D(416, 92), 210, 12, "anchor threaded into the next", 9.5, False, GRAY_MID)
    text(sh, *D(416, 105), 200, 12, "two rounds", 9.5, True, NAVY)

    # ═══ BAND 02 — team size ════════════════════════════════════════════
    by = BANDS[1] + DIAG_DY

    # MINIMAL
    text(sh, *D(0, -2), 260, 12, "M I N I M A L  —  A N S W E R - D R I V E N",
         10, True, BLUE)
    label_in(rect(sh, *D(0, 50), 40, 22, NAVY), "ORCH", 8, True, WHITE)
    arrow(sh, *D(40, 61), *D(62, 61))
    chip(sh, DX + 70, by + 50, 150, 22, "bureau")
    text(sh, *D(228, 50), 120, 22, "fetch a number", 8.5, False, GRAY_MID)
    text(sh, *D(0, 92), 240, 12, "every pick carries weight —", 9.5, False, GRAY_MID)
    text(sh, *D(0, 105), 240, 12, "no “in case relevant”", 9.5, True, NAVY)

    # MAXIMAL
    text(sh, *D(300, -2), 300, 12,
         "M A X I M A L  —  D I S C O V E R Y - D R I V E N", 10, True, BLUE)
    label_in(rect(sh, *D(300, 28), 40, 48, NAVY), "ORCH", 8, True, WHITE)
    for ey in (23, 44, 65, 86):
        arrow(sh, *D(340, 52), *D(358, ey), GRAY_MID, 0.9)
    fan = [(14, "modeling", "a score that lagged the deterioration", 440),
           (35, "spend_payments", "behaviour that moved first", 454),
           (56, "bureau", "did external call it earlier?", 440),
           (77, "crossbu", "exposure the scoring can’t see", 440)]
    for ry, nm, note, nx in fan:
        rect(sh, DX + 364, by + ry, 242, 17, BOX_BG, BOX_LN)
        text(sh, DX + 370, by + ry, 90, 17, nm, 8.5, True, NAVY)
        text(sh, DX + nx, by + ry, 170, 17, note, 8.5, False, GRAY_MID)
    text(sh, *D(300, 98), 300, 12, "4+ in parallel, one round —", 9.5, False, GRAY_MID)
    text(sh, *D(300, 111), 320, 12, "orthogonal directions, not duplicate votes",
         9.5, True, NAVY)

    # ═══ BAND 03 — coherence review ═════════════════════════════════════
    by = BANDS[2] + DIAG_DY

    text(sh, *D(0, -2), 300, 12,
         "F I R E S  O N  ≥ 2  S P E C I A L I S T S", 10, True, BLUE)
    chip(sh, DX + 0, by + 30, 104, 18, "spend_payments")
    chip(sh, DX + 0, by + 54, 104, 18, "modeling")
    arrow(sh, *D(104, 39), *D(130, 47))
    arrow(sh, *D(104, 63), *D(130, 55))

    label_in(rect(sh, *D(138, 30), 112, 42, NAVY),
             [("general_specialist", 8.5, True, WHITE),
              ("verify-only tools", 7.5, False, PALE_BLUE)])
    arrow(sh, *D(250, 51), *D(276, 51))

    rect(sh, DX + 284, by + 14, 218, 74, WHITE, BLUE)
    text(sh, *D(294, 20), 120, 12, "D I R E C T I V E", 8, True, BLUE)
    directive = [(36, GREEN, "coherent — ship it", GREEN_D, False),
                 (54, BLUE, "needs_redispatch — re-run ONE, anchored", NAVY, True),
                 (72, GRAY_MID, "qualified_release — ship with a flag", GRAY_TX, False)]
    for ry, sw, body, col, bold in directive:
        rect(sh, DX + 294, by + ry - 2, 5, 5, sw)
        text(sh, DX + 305, by + ry - 6, 200, 13, body, 8.5, bold, col)

    arrow(sh, *D(502, 51), *D(526, 51), BLUE, 1.25)
    label_in(rect(sh, *D(534, 40), 72, 22, BOX_BG, BOX_LN), "synthesize", 8.5, True)

    text(sh, *D(0, 98), 606, 12,
         "it re-measures, it never dispatches — the directive is advisory "
         "and the orchestrator acts on it", 9.5, False, GRAY_MID)
    text(sh, *D(0, 111), 606, 12,
         "capped at 2 rounds · any reviewer failure degrades to the answer "
         "we already had", 9.5, True, NAVY)

    # ── band dividers + footnote ────────────────────────────────────────
    for y in (278, 480):
        hline(sh, 44, y, 1192)
    hline(sh, 44, 682, 1192)
    text(sh, 44, 690, 1192, 16,
         "One question, three decisions: **what shape**, **how many**, and "
         "**who checks the result**.",
         11.5, False, GRAY_MID, align=PP_ALIGN.CENTER)

    # ── speaker notes ───────────────────────────────────────────────────
    slide.notes_slide.notes_text_frame.text = (
        "01 Dispatch shape - parallel is the default. Collapse is the one to "
        "call out: hand the whole causal chain to ONE cross-querying "
        "specialist that self-anchors, and the dependency costs no extra "
        "round. Sequential only when the anchor is itself heavy.\n\n"
        "02 Team size - the minimum-set rule is deliberately INVERTED for "
        "open, senior-level questions (no metric named, no table named). "
        "There, each specialist gets a direction of investigation rather than "
        "a fetch, and the directions are orthogonal so results compose "
        "instead of voting on one point.\n\n"
        "03 General specialist review - this one is server-enforced, not "
        "chosen by the orchestrator. It fires on every turn with 2+ domain "
        "specialists. It carries verify-only tools (it can re-measure an "
        "aggregate, it cannot dump rows) and it never dispatches - the "
        "directive is advisory. If it fails or times out, we ship the answer "
        "we already had.\n\n"
        "The 2025-05 / 2024-09 example runs through all three bands: it is "
        "what naive parallel gets wrong, what collapse/sequential prevents, "
        "and what the review catches when planning missed it."
    )


# ── slide 2 — distiller ─────────────────────────────────────────────────────
def slide_distiller(prs):
    """Condensed from brainstorm/layer1_memory_distiller.svg — the thesis of
    that diagram is that ONE knowledge point has TWO provenances, so the slide
    keeps the provenance colouring (blue = LLM, green = deterministic) and
    drops the full orchestrator/specialist cycle."""
    slide = new_slide(prs)
    sh = slide.shapes
    head(sh, "Distiller | how a finding becomes memory")

    # ═══ BAND 1 — where it sits (the fork) ══════════════════════════════
    text(sh, 44, 84, 300, 12, "W H E R E   I T   S I T S", 10, True, BLUE)

    label_in(rect(sh, 44, 106, 140, 30, BOX_BG, BOX_LN), "SpecialistOutput", 9)
    arrow(sh, 184, 121, 208, 108)
    arrow(sh, 184, 121, 208, 148)

    label_in(rect(sh, 214, 96, 286, 26, BOX_BG, BOX_LN),
             "synthesis → FinalAnswer → reviewer", 8.5)
    text(sh, 512, 96, 400, 26, "the answer path — **writes nothing to the KB**",
         9.5, False, GRAY_MID)

    label_in(rect(sh, 214, 134, 286, 26, NAVY), "Distiller → knowledge points", 8.5,
             True, WHITE)
    text(sh, 512, 134, 400, 26, "the memory path — **the only writer**",
         9.5, False, GRAY_MID)

    rect(sh, 920, 92, 316, 72, GRAY_BG)
    rect(sh, 920, 92, 4, 72, BLUE)
    text(sh, 936, 100, 292, 58,
         "Both run **in parallel**. A distiller failure costs the memory, "
         "never the answer.", 11.5, False, GRAY_TX,
         anchor=MSO_ANCHOR.MIDDLE, wrap=True, spacing=1.2)

    hline(sh, 44, 178, 1192)

    # ═══ BAND 2 — the assembly ══════════════════════════════════════════
    # ── column A: the two sources ──
    text(sh, 44, 192, 300, 12, "T W O   S O U R C E S", 10, True, BLUE)

    rect(sh, 44, 214, 256, 96, INK)
    text(sh, 52, 218, 200, 12, "TOOL OUTPUT · JSON", 8, True, PALE_BLUE)
    code(sh, 44, 236, 256, 72, [
        '"groups": [',
        ' {"group":"AMAZON",  "raw_value":45200},',
        ' {"group":"WALMART", "raw_value":32100},',
        ' {"group":"COSTCO",  "raw_value":28500} ],',
        '"concentration":{"hhi":0.27,"top3":0.80}',
    ], 7.5, CODE_TX)
    text(sh, 44, 316, 300, 12,
         "ground truth — **never re-enters an LLM**", 9.5, False, GREEN_D)

    rect(sh, 44, 340, 256, 62, BOX_BG, BOX_LN)
    text(sh, 52, 344, 200, 12, "FINDINGS + EVIDENCE", 8, True, BLUE)
    text(sh, 54, 360, 248, 40,
         "“Spend moderately concentrated: AMAZON leads $45,200 (34%), "
         "top-3 = 80% (HHI 0.27).”", 8.5, False, GRAY_TX,
         anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.15)
    text(sh, 44, 408, 300, 12, "the specialist's own prose", 9.5, False, GRAY_MID)

    arrow(sh, 302, 306, 328, 306, GRAY_MID, 1.2)

    # ── column B: the distiller's two lanes ──
    text(sh, 336, 192, 300, 12, "T W O   L A N E S", 10, True, BLUE)

    label_in(rect(sh, 336, 214, 324, 24, NAVY), "DISTILLER", 9, True, WHITE)

    rect(sh, 336, 244, 324, 74, WHITE, BLUE)
    text(sh, 344, 249, 200, 12, "① Agent skeleton — LLM", 8.5, True, BLUE)
    text(sh, 500, 249, 154, 12, "← from findings", 7.5, False, GRAY_MID,
         align=PP_ALIGN.RIGHT)
    code(sh, 336, 266, 324, 50, [
        ('"topic":   ‹LLM›', BLUE),
        ('"claim":   ‹LLM›', BLUE),
        ('"numbers": ∅  ← deterministic', GRAY_MID),
        ('"viz":     ∅  ← auto-chart', GRAY_MID),
    ], 7.5)

    rect(sh, 336, 326, 324, 74, WHITE, GREEN)
    text(sh, 344, 331, 220, 12, "② _ParsedSeries — no LLM", 8.5, True, GREEN_D)
    text(sh, 500, 331, 154, 12, "← from tool output", 7.5, False, GRAY_MID,
         align=PP_ALIGN.RIGHT)
    code(sh, 336, 348, 324, 50, [
        ('column_name = "amount_by_merchant"', GRAY_TX),
        ('key_field   = "group"', GRAY_TX),
        ('lookup = {"AMAZON":45200, "WALMART":32100,', GREEN_D),
        ('          "COSTCO":28500, ...}', GREEN_D),
    ], 7.5)

    text(sh, 336, 408, 320, 12, "summarize_by_group → key + raw_value",
         9.5, False, GRAY_MID)

    arrow(sh, 662, 306, 688, 306, GRAY_MID, 1.2)

    # ── column C: the knowledge point ──
    text(sh, 696, 192, 400, 12, "O N E   K N O W L E D G E   P O I N T",
         10, True, BLUE)

    label_in(rect(sh, 696, 214, 540, 24, NAVY),
             "③  MERGED — written to the session KB", 9, True, WHITE)

    kp_rows = [
        (244, 30, BLUE, "topic", '"spend_by_merchant"', "LLM", BLUE),
        (278, 44, BLUE, "claim",
         '"AMAZON leads $45,200 (34%);\ntop-3 = 80% (HHI 0.27)"', "LLM", BLUE),
        (326, 44, GREEN, "numbers",
         '[{"group":"AMAZON","value":45200},\n {"group":"WALMART","value":32100}, …]',
         "DETERMINISTIC", GREEN_D),
        (374, 30, GREEN, "viz", '/charts/spend_by_merchant_share.png',
         "BY PATH", GREEN_D),
    ]
    for y, h, bar, field, value, tag, tag_col in kp_rows:
        rect(sh, 696, y, 540, h, WHITE, GRAY_LN)
        rect(sh, 696, y, 4, h, bar)
        text(sh, 712, y, 70, h, field, 9, True, NAVY)
        code(sh, 786, y + 6, 350, h - 6,
             [(ln, GRAY_TX) for ln in value.split("\n")], 7.5)
        text(sh, 1120, y, 108, h, tag, 7.5, True, tag_col, align=PP_ALIGN.RIGHT)

    text(sh, 696, 412, 540, 12,
         "one object, two provenances — **the LLM never writes a number**",
         9.5, False, GRAY_MID)

    hline(sh, 44, 438, 1192)

    # ═══ BAND 3 — three takeaways ═══════════════════════════════════════
    cards = [
        ("Why parse, not prompt?",
         "**a · the firewall** — raw numbers never re-enter an LLM, so there is "
         "nothing to redact and nothing to reject.\n"
         "**b · faithfulness** — KB values exactly equal tool output. No drift, "
         "no hallucinated figure."),
        ("Auto-chart — no model call",
         "Validate the tool-output schema → inject threshold lines from the data "
         "catalog → stream the spec over SSE → store the image **path**, not "
         "pixels. Deterministic end to end."),
        ("What the next turn gets",
         "A follow-up reads the **exact numbers**, not a paraphrase of them — "
         "which is what makes the cross-case layer safe to build on."),
    ]
    for i, (title, body) in enumerate(cards):
        x = 44 + i * 412
        rect(sh, x, 456, 368, 158, GRAY_BG)
        rect(sh, x, 456, 4, 158, BLUE)
        text(sh, x + 18, 472, 330, 22, title, 15, True, NAVY,
             anchor=MSO_ANCHOR.TOP, wrap=True)
        text(sh, x + 18, 502, 332, 130, body, 11.5, False, GRAY_TX,
             anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.28)

    hline(sh, 44, 664, 1192)
    text(sh, 44, 674, 1192, 18,
         "The LLM writes **what it means**. The parser writes "
         "**what it measured**.",
         11.5, False, GRAY_MID, align=PP_ALIGN.CENTER)

    slide.notes_slide.notes_text_frame.text = (
        "The point of this slide is provenance, not plumbing.\n\n"
        "A knowledge point has four fields. Two are written by the LLM (topic, "
        "claim) and two are filled deterministically (numbers from the parsed "
        "tool output, viz as a file path from the auto-chart pipeline). The "
        "LLM never writes a number into memory.\n\n"
        "Why it is built that way: (a) the firewall - raw numbers never "
        "re-enter an LLM, so there is nothing to redact; (b) faithfulness - "
        "KB values exactly equal tool output, so a follow-up turn reads the "
        "real figure rather than a paraphrase that drifted.\n\n"
        "_ParsedSeries is the midterm product that proves the path: it is "
        "built by parsing summarize_by_group output into key + raw_value, with "
        "no model in the loop.\n\n"
        "The distiller runs in parallel to the answer path. FinalAnswer streams "
        "to the reviewer and writes nothing to the KB, so a distiller failure "
        "costs the memory, never the answer.\n\n"
        "Full engineering diagram: brainstorm/layer1_memory_distiller.svg"
    )


# ── slide 4 — stall-and-retry vs self-recovery, in detail ──────────────────────────────
def slide_recovery_detail(prs):
    """Two recovery machines that are constantly confused for each other.
    Sources: brainstorm/timeouts-and-retries.md (fences, measured stalls, the
    multiplication trap) and agent_factories/agent_tools/agent_tool.py:590-660
    (the three self-recovery triggers)."""
    slide = new_slide(prs)
    sh = slide.shapes
    head(sh, "Recovery | the mechanism")

    # ═══ LEFT — stall-and-retry ═════════════════════════════════════════
    text(sh, 44, 84, 400, 14, "S T A L L - A N D - R E T R Y", 11, True, BLUE)
    text(sh, 44, 102, 400, 14, "did anything come back?   —   a clock decides",
         11.5, False, GRAY_TX)

    ladder = [
        (44, 572, "TURN", "360s · cancelled — every agent in flight dies with it",
         "no retry"),
        (60, 556, "PHASE", "screen 30s · orch_plan 25s · reviewer 120s · specialist 240s",
         "2 retry"),
        (76, 540, "AGENT", "the whole agentic loop re-runs", "2 attempts"),
        (92, 524, "CALL", "stall fence 40s → abandon and re-issue", "once"),
    ]
    for i, (x, w, name, detail, badge) in enumerate(ladder):
        y = 126 + i * 38
        rect(sh, x, y, w, 32, BOX_BG, BOX_LN)
        rect(sh, x, y, 4, 32, NAVY if i < 2 else BLUE)
        text(sh, x + 14, y, 78, 32, name, 9.5, True, NAVY)
        text(sh, x + 96, y, w - 190, 32, detail, 8.5, False, GRAY_TX)
        text(sh, x + w - 86, y, 76, 32, badge, 8.5, True, BLUE,
             align=PP_ALIGN.RIGHT)

    text(sh, 44, 282, 572, 14,
         "phase expiry = that agent failed; **the turn continues, degraded**",
         9.5, False, GRAY_MID)

    rect(sh, 44, 304, 572, 96, GRAY_BG)
    rect(sh, 44, 304, 4, 96, BLUE)
    text(sh, 60, 312, 540, 14,
         "measured: safechain calls do not run slow — **they stall**",
         10.5, True, NAVY)
    code(sh, 60, 332, 540, 46, [
        ("report_agent         100.01s   ← its fence was 100", GRAY_TX),
        ("distiller.modeling   120.01s   ← distiller_s = 120", GRAY_TX),
        ("general_specialist    25.00s   ← orch_plan_s = 25", GRAY_TX),
    ], 8)
    text(sh, 60, 380, 540, 14,
         "every failure died at **its own fence**, not at a natural duration — "
         "normal calls in those same turns took 2–13s", 8.5, False, GRAY_MID,
         wrap=True)

    # ═══ centre divider ═════════════════════════════════════════════════
    cn = sh.add_connector(MSO_CONNECTOR.STRAIGHT, px(640), px(100), px(640), px(400))
    cn.line.color.rgb = RGBColor(0xC7, 0xCF, 0xDA)
    cn.line.width = Pt(1)

    # ═══ RIGHT — self-recovery ══════════════════════════════════════════
    text(sh, 664, 84, 400, 14, "S E L F - R E C O V E R Y", 11, True, GREEN_D)
    text(sh, 664, 102, 500, 14, "is what came back usable?   —   the content decides",
         11.5, False, GRAY_TX)

    triggers = [
        ("ungrounded_retry", "a tool call under it failed",
         "the tool errors are fed back into the re-run"),
        ("absence_reread", "the answer denies rows the tool returned",
         "re-reads the transcript it has — no new query"),
        ("max_turns_retry", "blew its turn budget",
         "collapses to a single decisive call"),
        ("guidance loop", "403 / 400 firewall rejection",
         "appends what to remove, then retries"),
    ]
    for i, (name, detects, changes) in enumerate(triggers):
        y = 126 + i * 44
        rect(sh, 664, y, 572, 38, WHITE, GRAY_LN)
        rect(sh, 664, y, 4, 38, GREEN)
        code(sh, 676, y + 5, 150, 14, [(name, GREEN_D)], 8.5)
        text(sh, 676, y + 20, 200, 13, detects, 8, False, GRAY_MID)
        text(sh, 880, y, 348, 38, changes, 8.5, False, GRAY_TX, wrap=True)

    rect(sh, 664, 314, 572, 86, GRAY_BG)
    rect(sh, 664, 314, 4, 86, GREEN)
    text(sh, 680, 322, 540, 14, "the rule that separates them", 10.5, True, NAVY)
    text(sh, 680, 342, 542, 52,
         "**401 is transient** — a stale token, so it refreshes and retries "
         "**in place**. **403 / 400 are deterministic** — an identical retry is "
         "guaranteed to fail the same way, so the input itself must change.",
         9.5, False, GRAY_TX, anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.25)

    hline(sh, 44, 424, 1192)

    # ═══ BOTTOM — the three traps ═══════════════════════════════════════
    traps = [
        ("The multiplication trap",
         "Agent-level and call-level retries **compose**. A specialist that "
         "retries once, each of whose calls stalls once, issues **up to 4 "
         "requests** — and one phase fence holds all of them. That is why the "
         "call budget is exactly one extra attempt."),
        ("A fence can make a retry unreachable",
         "**screen** at 30s sat **under** the 40s call-layer stall fence, so that "
         "retry was dead code for the phase. Fixed by giving screen its own "
         "**10s fence** — not by loosening the phase."),
        ("An abandoned attempt still costs",
         "If inference had started, the tokens are **billed** — cancelling "
         "closes your side, it does not un-run the work. node_trace records "
         "**zero** for it, so trace totals are a floor, not a total."),
    ]
    for i, (title, body) in enumerate(traps):
        x = 44 + i * 412
        rect(sh, x, 446, 368, 172, GRAY_BG)
        rect(sh, x, 446, 4, 172, BLUE)
        text(sh, x + 18, 462, 330, 20, title, 14, True, NAVY,
             anchor=MSO_ANCHOR.TOP, wrap=True)
        text(sh, x + 18, 490, 332, 118, body, 11, False, GRAY_TX,
             anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.28)

    hline(sh, 44, 646, 1192)
    text(sh, 44, 656, 1192, 18,
         "Widening a fence only converts a failure into a 130-second wait. "
         "**Escaping the stall is what recovers the time.**",
         11.5, False, GRAY_MID, align=PP_ALIGN.CENTER)

    slide.notes_slide.notes_text_frame.text = (
        "The one thing to land: these are two different machines answering two "
        "different questions. Stall-and-retry asks 'did anything come back' and "
        "a clock decides. Self-recovery asks 'is what came back usable' and the "
        "content decides.\n\n"
        "Left - four nested fences. Turn 360s is a hard cancel with no retry. "
        "Phase expiry means that agent failed but the turn continues degraded. "
        "Only two phases retry: orch_plan replans, screen re-issues. The call "
        "layer abandons at a 40s stall and re-issues exactly once.\n\n"
        "The evidence for building it this way: safechain calls do not run "
        "slow, they stall. Failures died at exactly their own fence - 100.01, "
        "120.01, 25.00 - while normal calls in the same turns took 2 to 13 "
        "seconds. A number that lands on the fence to two decimals is a stall, "
        "not slow work.\n\n"
        "Right - self-recovery retries with a CHANGED input, because the answer "
        "was wrong rather than late. absence_reread is the interesting one: "
        "nothing broke, the specialist misread a correct result (one returned "
        "payment in, 'zero records' out), so the retry is a re-read of the "
        "transcript it already has, not a fresh query.\n\n"
        "If asked why not just retry everything: the multiplication trap. Two "
        "agent attempts times two call attempts is four requests under one "
        "phase fence, and the arithmetic has to stay inspectable."
    )


# ── slide 3 — recovery, the concept ─────────────────────────────────────────
def slide_recovery_concept(prs):
    """Concept first: one fork, two machines, four short lines. Every number
    and mechanism lives on the following slide."""
    slide = new_slide(prs)
    sh = slide.shapes
    head(sh, "Recovery | two failures, two machines")

    # ── the fork ────────────────────────────────────────────────────────
    label_in(rect(sh, 500, 96, 280, 44, NAVY), "a turn goes wrong", 15, True, WHITE)
    arrow(sh, 596, 142, 400, 186, GRAY_MID, 1.4)
    arrow(sh, 684, 142, 880, 186, GRAY_MID, 1.4)

    # ── the two panels ──────────────────────────────────────────────────
    panels = [
        (100, BLUE, "NOTHING CAME BACK", "Stall-and-retry",
         "A clock decides.", "Retry the SAME request.",
         "The call goes quiet — abandon it and re-issue, "
         "instead of waiting out the budget."),
        (708, GREEN, "SOMETHING CAME BACK — AND IT'S WRONG", "Self-recovery",
         "The content decides.", "Retry with a CHANGED input.",
         "The answer denies rows the tool returned — re-read what it "
         "already has, don't re-query."),
    ]
    for x, accent, eyebrow, name, decides, move, example in panels:
        rect(sh, x, 190, 472, 356, WHITE, GRAY_LN)
        rect(sh, x, 190, 5, 356, accent)
        text(sh, x + 32, 214, 420, 16, eyebrow, 10.5, True, accent)
        text(sh, x + 32, 240, 420, 40, name, 28, True, NAVY)
        hline(sh, x + 32, 292, 96, accent, 2.0)
        text(sh, x + 32, 306, 420, 24, decides, 15, False, GRAY_TX)
        text(sh, x + 32, 336, 420, 24, move, 15, True, NAVY)
        rect(sh, x + 32, 390, 408, 124, GRAY_BG)
        text(sh, x + 50, 412, 376, 84, example, 12.5, False, GRAY_TX,
             anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.3)

    # ── the one takeaway ────────────────────────────────────────────────
    hline(sh, 44, 584, 1192)
    text(sh, 44, 608, 1192, 30,
         "Same word — **retry**.  Opposite mechanics: "
         "**one repeats the request, the other changes it.**",
         17, False, GRAY_TX, align=PP_ALIGN.CENTER)
    text(sh, 44, 652, 1192, 20,
         "An identical retry only helps when the failure was transient.",
         12, False, GRAY_MID, align=PP_ALIGN.CENTER)

    slide.notes_slide.notes_text_frame.text = (
        "Lead with the fork, not the fences.\n\n"
        "Everything downstream follows from one question: did anything come "
        "back at all? If nothing came back, no amount of reasoning helps - you "
        "abandon and re-issue the same request, and a clock is the only thing "
        "that can make that call. If something DID come back and it is wrong, "
        "repeating the identical request is guaranteed to reproduce the same "
        "wrong result, so the input has to change.\n\n"
        "That is the whole concept. The next slide is the mechanism - fences, "
        "triggers, and the traps - and it is there for the questions, not for "
        "the walkthrough."
    )


# ── slide 5 — self-recovery, worked examples ────────────────────────────────
def slide_recovery_examples(prs):
    """The `_DEGRADED_RECOVERY` map (agent_tool.py:116) verbatim in spirit:
    every classified tool failure is paired with the specific next move, because
    the specialist gets exactly one retry to use it."""
    slide = new_slide(prs)
    sh = slide.shapes
    head(sh, "Self-recovery | the failure carries its own fix")

    # ── the loop, in one strip ──────────────────────────────────────────
    strip = [
        (44, "tool call fails", RED),
        (288, "classified by marker", BLUE),
        (532, "repair line appended", BLUE),
        (776, "ONE retry", NAVY),
        (1020, "grounded answer", GREEN),
    ]
    for i, (x, lbl, col) in enumerate(strip):
        filled = col is NAVY
        box = rect(sh, x, 92, 192, 32,
                   NAVY if filled else BOX_BG, None if filled else col)
        label_in(box, lbl, 9.5, True, WHITE if filled else col)
        if i < len(strip) - 1:
            arrow(sh, x + 196, 108, x + 240, 108, GRAY_MID, 1.1)

    # ── the examples table ──────────────────────────────────────────────
    COLS = [(44, 212), (272, 330), (618, 618)]
    hdrs = ["WHAT WENT WRONG", "WHAT ACTUALLY HAPPENED",
            "THE REPAIR HANDED BACK — one retry, so it must be actionable"]
    for (x, w), t in zip(COLS, hdrs):
        text(sh, x, 152, w, 14, t, 9, True, BLUE)
    hline(sh, 44, 172, 1192, GRAY_LN, 1.0)

    rows = [
        ("column_not_found", "named a column that table does not carry",
         "call **search_columns** or **get_table_schema**, then re-issue with a real "
         "column name — and do **NOT** report it as a data gap", False),
        ("no_buckets", "the date column's format did not parse — “Jul-25” "
         "where ISO was assumed",
         "read the column's **actual format** via get_table_schema and re-issue "
         "with the correct date_column", False),
        ("specs_unparseable", "malformed specs_json — **the tool never ran at all**",
         "re-send as a valid JSON array, or fall back to the single-spec tool "
         "**one call at a time**", False),
        ("table_not_found", "that table does not exist for this case",
         "call **list_available_tables** and pick one that does", False),
        ("(nothing failed)", "the filter matched **0 rows** — and an empty result "
         "looks exactly like a true negative",
         "**_why_empty** re-runs each condition alone: “**NOT A NEGATIVE FINDING** — "
         "these conditions name columns this table does not have, so they were "
         "never tested”", True),
    ]
    y = 182
    for name, happened, repair, star in rows:
        h = 62 if star else 54
        if star:
            rect(sh, 44, y, 1192, h, GRAY_BG)
            rect(sh, 44, y, 4, h, RED)
        code(sh, COLS[0][0] + (14 if star else 0), y + 8, COLS[0][1], 16,
             [(name, RED if star else NAVY)], 9)
        text(sh, COLS[1][0], y, COLS[1][1], h, happened, 10, False, GRAY_TX,
             wrap=True, spacing=1.2)
        text(sh, COLS[2][0], y, COLS[2][1], h, repair, 10, False,
             NAVY if star else GRAY_TX, wrap=True, spacing=1.2)
        y += h + 4
        if not star:
            hline(sh, 44, y - 2, 1192, RGBColor(0xEE, 0xF1, 0xF5), 0.75)

    # ── the three principles ────────────────────────────────────────────
    hline(sh, 44, 512, 1192)
    notes = [
        ("One retry means the message must carry the next move.",
         "“your call failed” is not enough — every line names the specific tool "
         "to call and what to change."),
        ("The dangerous failure is the one that does not fail.",
         "0 rows raises nothing. It has to be made to speak, or it is reported "
         "as a finding."),
        ("Never let a broken call become a finding.",
         "each repair line ends by forbidding the false report — “do NOT report "
         "it as a data gap”."),
    ]
    for i, (t, b) in enumerate(notes):
        x = 44 + i * 412
        rect(sh, x, 532, 368, 116, GRAY_BG)
        rect(sh, x, 532, 4, 116, BLUE)
        text(sh, x + 18, 546, 332, 34, t, 12.5, True, NAVY,
             anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.18)
        text(sh, x + 18, 588, 332, 52, b, 10.5, False, GRAY_TX,
             anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.22)

    text(sh, 44, 664, 1192, 20,
         "A failed tool call is not a dead end — it is an instruction.",
         12, False, GRAY_MID, align=PP_ALIGN.CENTER)

    slide.notes_slide.notes_text_frame.text = (
        "These are the real classifications from _DEGRADED_RECOVERY in "
        "agent_tool.py - each key is a tool failure the grounding layer can "
        "detect, each value is the guidance appended to the retry.\n\n"
        "The design rule is in the code comment: keep each line ACTIONABLE, "
        "because the model gets one retry, so 'your call failed' is not enough "
        "- it needs the specific next move.\n\n"
        "The bottom row is the one worth dwelling on. Nothing failed. The tool "
        "succeeded and returned zero rows, which is indistinguishable between "
        "a wrong column name, a wrong value vocabulary, and a genuine absence. "
        "The observed failure: a specialist reported no high-TSR transactions "
        "on a case that had them, because the filter named a score column on a "
        "table that only carries amount and merchant. That condition was never "
        "tested. _why_empty now re-runs each condition alone and leads with "
        "'NOT A NEGATIVE FINDING' when a named column is missing.\n\n"
        "Note the guard clauses - 'do NOT report it as a data gap', 'do NOT "
        "report this as no such rows exist'. Half the value here is stopping a "
        "broken call from being laundered into a confident negative."
    )


# ── slide 6 — grounded rate on a discovery question ─────────────────────────
def slide_grounded_rate(prs):
    """“Any model opportunities?” scored against the REAL specialist answer.
    The 20.9% row is a genuine defect in that answer: its own cited endpoints
    (787 → 606) give 23.0%. NOTE: the metric is computed in AgenticEval's
    content/ pipeline, not in this repo."""
    slide = new_slide(prs)
    sh = slide.shapes
    head(sh, "Grounded rate | scoring “any model opportunities?”")

    # ══ the raw answer ══════════════════════════════════════════════════
    text(sh, 44, 82, 500, 13, "T H E   R A W   A N S W E R", 9.5, True, BLUE)
    rect(sh, 44, 100, 1192, 104, BOX_BG, BOX_LN)
    text(sh, 62, 108, 1160, 90,
         "“No prior curated reports — answer is from live specialist analysis "
         "only.  Internal risk scores (CDSS, TSR) reacted late: spend-"
         "concentration breach (>2.4, peak 8.73) occurred Dec'2025–Jan'2026, but "
         "model scores escalated only from Feb'2026, **missing actionable "
         "windows**.  All spend (**$44,912.26**) went to OMEGA 3 MOTOR SPORTS in "
         "Jan'2026.  Model risk-driving features **were redundant**, dominated by "
         "'ioip_interaction' across 10+ months.  Card exposure is totally driven "
         "by BUSINESS GOLD REWARDS (100% of balance, 93% of limit); BLUE BUSINESS "
         "PLUS has $700 unused headroom at zero balance, **representing "
         "unrealized mitigation avenues**.  Bureau FICO dropped sharply (787 → 606 "
         "in May'2026, **a 20.9% decline**) without integration as a model driver, "
         "**indicating missed triggers for early warning**.”",
         10, False, GRAY_TX, anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.28)

    # ══ the claims ══════════════════════════════════════════════════════
    for x, w, t in ((44, 470, "A T O M I C   C L A I M"),
                    (528, 118, "V E R D I C T"),
                    (664, 572, "T H E   R E A S O N I N G")):
        text(sh, x, 214, w, 13, t, 9.5, True, BLUE)
    text(sh, 900, 214, 336, 13, "8 of 12 eligible claims shown", 8.5, False,
         GRAY_MID, align=PP_ALIGN.RIGHT)
    hline(sh, 44, 232, 1192, GRAY_LN, 1.0)

    claims = [
        ("All spend — $44,912.26 — went to one merchant, Jan'26", "supported",
         "summarize_by_group(spends, by merchant): a single group, top1 = 1.00",
         GREEN_D),
        ("Concentration breached >2.4, peaking at 8.73, Dec'25–Jan'26",
         "supported", "aggregate_column(op=max) over that window returns 8.73",
         GREEN_D),
        ("Model scores escalate only from Feb'26", "supported",
         "summarize_trend per month — first escalation month is 2026-02", GREEN_D),
        ("→ risk scoring lags spend by ~1–2 months", "supported",
         "derived: the gap between two measured series, both in the trace",
         GREEN_D),
        ("FICO 787 (Aug'25) → 606 (May'26), no new delinquencies", "supported",
         "bureau trend + delinquent_external_trades = 0 over the same window",
         GREEN_D),
        ("No payment returns; no bureau-inquiry data exists", "supported",
         "both are **checked non-findings** — a checked non-finding is a result",
         GREEN_D),
        ("“…a 20.9% decline”", "UNSUPPORTED",
         "its own cited endpoints give **23.0%**. The derived value does not "
         "follow from the evidence beside it.", RED),
        ("“the drivers were redundant”", "UNSUPPORTED",
         "redundancy IS countable — driver correlation — but the trace only "
         "ranked drivers by frequency. Never computed.", RED),
    ]
    for i, (claim, verdict, why, col) in enumerate(claims):
        y = 240 + i * 29
        bad = col is RED
        if bad:
            rect(sh, 44, y, 1192, 27, GRAY_BG)
        rect(sh, 44, y, 3, 27, col)
        text(sh, 60, y, 460, 27, claim, 9.5, bad, NAVY if bad else GRAY_TX)
        text(sh, 528, y, 118, 27, verdict, 9, True, col)
        text(sh, 664, y, 572, 27, why, 9, False, GRAY_MID, wrap=True, spacing=1.08)

    # ══ ineligible — three different reasons ════════════════════════════
    text(sh, 44, 484, 700, 13,
         "N O T   E L I G I B L E   —   A N D   W H Y", 9.5, True, GRAY_MID)
    rect(sh, 44, 502, 700, 148, GRAY_BG)
    rect(sh, 44, 502, 4, 148, GRAY_MID)
    inelig = [
        ("“missing actionable windows”", "NO REFERENT",
         "no column defines an actionable window — the lag is countable, "
         "the loss is not"),
        ("“unrealized mitigation avenues”", "JUDGEMENT",
         "an interpretation laid over a supported fact; nothing marks headroom "
         "as opportunity"),
        ("“missed triggers for early warning”", "COUNTERFACTUAL",
         "asserts what would have happened — the data records what did"),
    ]
    for i, (phrase, kind, why) in enumerate(inelig):
        y = 512 + i * 45
        text(sh, 64, y, 290, 15, phrase, 10, True, NAVY)
        text(sh, 64, y + 16, 290, 13, kind, 8, True, GRAY_MID)
        text(sh, 372, y, 356, 32, why, 9.5, False, GRAY_TX, wrap=True,
             spacing=1.12)

    # ══ the two numbers ═════════════════════════════════════════════════
    rect(sh, 768, 502, 226, 148, WHITE, BLUE)
    rect(sh, 768, 502, 4, 148, BLUE)
    text(sh, 790, 516, 190, 13, "GROUNDED RATE", 10, True, BLUE)
    text(sh, 790, 534, 190, 42, "10 / 12", 30, True, NAVY)
    text(sh, 790, 582, 190, 17, "83%", 13, True, NAVY)
    text(sh, 790, 606, 194, 32, "supported ÷ eligible", 10, False, GRAY_TX,
         wrap=True)

    rect(sh, 1010, 502, 226, 148, WHITE, GREEN)
    rect(sh, 1010, 502, 4, 148, GREEN)
    text(sh, 1032, 516, 190, 13, "EVIDENCE LEVERAGE", 10, True, GREEN_D)
    text(sh, 1032, 534, 190, 42, "2.5", 30, True, NAVY)
    text(sh, 1032, 582, 196, 17, "claims per specialist", 11, True, NAVY)
    text(sh, 1032, 606, 198, 32,
         "4 dispatched — **read within question class**", 10, False, GRAY_TX,
         wrap=True)

    hline(sh, 44, 664, 1192)
    text(sh, 44, 676, 1192, 20,
         "Both failures are invisible to a reader: one number that does not "
         "follow, one claim that was never measured.",
         11.5, False, GRAY_MID, align=PP_ALIGN.CENTER)

    slide.notes_slide.notes_text_frame.text = (
        "This is the real answer to 'any model opportunities', scored claim by "
        "claim. Ten of twelve eligible claims are supported straight off the "
        "trace. Two are not, and they fail in completely different ways.\n\n"
        "The 20.9% is the one to stop on. The answer cites 787 in August 2025 "
        "and 606 in May 2026, and calls it a 20.9% decline. Those endpoints "
        "give 23.0%. Twenty point nine would require a starting score of 766. "
        "So the derived value does not follow from the evidence printed "
        "directly beside it - and note that every INPUT to it was correct. A "
        "human reading this answer would not catch that, because the numbers "
        "look cited and the sentence reads fluently.\n\n"
        "The second failure is different in kind: 'the drivers were redundant' "
        "is not a judgement. Redundancy is perfectly countable as a correlation "
        "between drivers. It fails because nobody computed it - the trace only "
        "ranked drivers by frequency. Unsupported, not ineligible, and that "
        "distinction is the whole discipline.\n\n"
        "The three ineligible lines fail three more ways: no referent, "
        "judgement, counterfactual. None of them is a defect - they are the "
        "interpretive layer an opportunity question is supposed to have. They "
        "just cannot be scored, so they leave the denominator.\n\n"
        "On leverage: a maximal team over-collects by design, so this number is "
        "only meaningful against other discovery turns. Grounded rate is the "
        "one that travels."
    )


def build(out_path: pathlib.Path) -> pathlib.Path:
    prs = Presentation()
    prs.slide_width = px(1280)
    prs.slide_height = px(720)
    slide_team_construction(prs)
    slide_distiller(prs)
    slide_recovery_concept(prs)
    slide_recovery_detail(prs)
    slide_recovery_examples(prs)
    slide_grounded_rate(prs)
    prs.save(str(out_path))
    return out_path


if __name__ == "__main__":
    out = pathlib.Path(__file__).with_name("2026-08-18-team-construction-slide.pptx")
    print("wrote", build(out))
